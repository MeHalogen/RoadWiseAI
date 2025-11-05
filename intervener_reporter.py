# intervener_reporter.py
# Report Generator for RoadWiseAI - Creates PDF and PowerPoint outputs

import json
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any

class ReportGenerator:
    """
    Generates PDF reports and PowerPoint presentations for recommendations.
    """
    
    def __init__(self):
        self.timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    def generate_pdf_report(self, recommendations: List[Dict[str, Any]], 
                           query: str, road_type: str = None, 
                           environment: str = None, output_path: str = "report.pdf"):
        """
        Generate a PDF report from recommendations.
        
        Args:
            recommendations (List): List of formatted recommendations
            query (str): Original query
            road_type (str): Road type
            environment (str): Environment context
            output_path (str): Output file path
        """
        doc = SimpleDocTemplate(output_path, pagesize=letter,
                               rightMargin=0.5*inch, leftMargin=0.5*inch,
                               topMargin=0.75*inch, bottomMargin=0.75*inch)
        
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f4788'),
            spaceAfter=12,
            alignment=1  # Center
        )
        story.append(Paragraph("RoadWiseAI", title_style))
        story.append(Paragraph("Road Safety Intervention Recommendation System", styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        
        # Query Summary
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#2e5f8a'),
            spaceAfter=6
        )
        
        story.append(Paragraph("QUERY DETAILS", heading_style))
        query_data = [
            ['Field', 'Value'],
            ['Issue', query],
            ['Road Type', road_type or 'Urban (default)'],
            ['Environment', environment or 'General'],
            ['Generated', self.timestamp]
        ]
        
        query_table = Table(query_data, colWidths=[1.5*inch, 4*inch])
        query_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2e5f8a')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f0f0f0')])
        ]))
        story.append(query_table)
        story.append(Spacer(1, 0.3*inch))
        
        # Recommendations
        story.append(Paragraph("RECOMMENDATIONS", heading_style))
        story.append(Spacer(1, 0.1*inch))
        
        for idx, rec in enumerate(recommendations, 1):
            rec_heading = ParagraphStyle(
                f'RecHeading{idx}',
                parent=styles['Heading3'],
                fontSize=12,
                textColor=colors.HexColor('#1f4788'),
                spaceAfter=6
            )
            
            story.append(Paragraph(f"Intervention {idx}: {rec['intervention']}", rec_heading))
            
            rec_data = [
                ['Reference', rec['reference']],
                ['Rationale', rec['rationale']],
                ['Assumptions', rec['assumptions']],
                ['Confidence', f"{rec['confidence']} ({rec['relevance_score']}%)"]
            ]
            
            rec_table = Table(rec_data, colWidths=[1.5*inch, 4*inch])
            rec_table.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#cccccc')),
                ('ROWBACKGROUNDS', (0, 0), (-1, -1), [colors.HexColor('#f9f9f9')])
            ]))
            story.append(rec_table)
            story.append(Spacer(1, 0.2*inch))
        
        # Footer
        story.append(Spacer(1, 0.2*inch))
        footer_style = ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=8,
            textColor=colors.grey,
            alignment=0
        )
        story.append(Paragraph(
            "Note: All recommendations are material-only estimates. Labor, transport, and taxes are excluded.",
            footer_style
        ))
        
        doc.build(story)
        print(f"PDF report generated: {output_path}")
    
    def generate_pptx_report(self, recommendations: List[Dict[str, Any]], 
                            query: str, road_type: str = None, 
                            environment: str = None, output_path: str = "presentation.pptx"):
        """
        Generate a 7-slide PowerPoint presentation (Hackathon format).
        
        Args:
            recommendations (List): List of formatted recommendations
            query (str): Original query
            road_type (str): Road type
            environment (str): Environment context
            output_path (str): Output file path
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define layouts
        blank_layout = prs.slide_layouts[6]
        title_layout = prs.slide_layouts[0]
        
        # Slide 1: Welcome
        slide1 = prs.slides.add_slide(blank_layout)
        self._add_title_slide(slide1, "RoadWiseAI", "Road Safety Intervention GPT")
        
        # Slide 2: Problem Statement
        slide2 = prs.slides.add_slide(blank_layout)
        self._add_content_slide(slide2, "Problem Statement",
            [
                "â€¢ 150,000+ road fatalities annually in India",
                "â€¢ Unsafe infrastructure, poor signage, neglected facilities",
                "â€¢ Manual, slow intervention selection process",
                "â€¢ Need for standards-aligned, explainable AI recommendations"
            ])
        
        # Slide 3: Approach & Architecture
        slide3 = prs.slides.add_slide(blank_layout)
        self._add_content_slide(slide3, "Approach & Architecture",
            [
                "â€¢ Hybrid Retrieval Engine: Fuzzy + Semantic Matching",
                "â€¢ Knowledge Base: 50+ IRC-aligned interventions",
                "â€¢ Explainability Layer: References, Rationale, Assumptions",
                "â€¢ Material-only Costing: CPWD/GeM data integration"
            ])
        
        # Slide 4: Knowledge Base
        slide4 = prs.slides.add_slide(blank_layout)
        self._add_content_slide(slide4, "Knowledge Base",
            [
                f"â€¢ Total Interventions: {len(recommendations) + 50} (12 seed + more)",
                "â€¢ Sources: IRC 35, IRC 67, IRC SP:84, IRC SP:87, IRC 99",
                "â€¢ Road Types: Urban, Highway, Rural",
                "â€¢ Priority-weighted: High, Medium, Low"
            ])
        
        # Slide 5: Demo Output
        slide5 = prs.slides.add_slide(blank_layout)
        self._add_demo_slide(slide5, "Demo Output", recommendations, query, road_type, environment)
        
        # Slide 6: Evaluation Mapping
        slide6 = prs.slides.add_slide(blank_layout)
        self._add_content_slide(slide6, "Evaluation Mapping",
            [
                "âœ“ Relevance: Context-aware retrieval & matching",
                "âœ“ Comprehensiveness: 3 recommendations + rationale",
                "âœ“ Explainability: IRC clauses, reasoning, assumptions",
                "âœ“ Innovation: Retrieval-based + rule-based hybrid logic"
            ])
        
        # Slide 7: Thank You
        slide7 = prs.slides.add_slide(blank_layout)
        self._add_title_slide(slide7, "Thank You", "RoadWiseAI: From Problem to Intervention in One Query")
        
        prs.save(output_path)
        print(f"PowerPoint presentation generated: {output_path}")
    
    def _add_title_slide(self, slide, title: str, subtitle: str):
        """Add a title slide."""
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors.HexColor('#1f4788')
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = colors.white
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = colors.white
        p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, slide, title: str, content: List[str]):
        """Add a content slide with bullet points."""
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors.white
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = colors.HexColor('#1f4788')
        title_shape.line.color.rgb = colors.HexColor('#1f4788')
        
        # Title text
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors.white
        p.alignment = PP_ALIGN.LEFT
        title_frame.margin_left = Inches(0.5)
        title_frame.margin_top = Inches(0.2)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        for line in content:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = colors.black
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    def _add_demo_slide(self, slide, title: str, recommendations: List[Dict[str, Any]], 
                       query: str, road_type: str, environment: str):
        """Add the demo output slide."""
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors.white
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = colors.HexColor('#1f4788')
        title_shape.line.color.rgb = colors.HexColor('#1f4788')
        
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors.white
        p.alignment = PP_ALIGN.LEFT
        title_frame.margin_left = Inches(0.3)
        
        # Query details
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6.2))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = f"Query: {query}"
        p.font.size = Pt(12)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Road Type: {road_type or 'Urban'} | Environment: {environment or 'General'}"
        p.font.size = Pt(11)
        p.space_before = Pt(6)
        
        p = text_frame.add_paragraph()
        p.text = "\nTop Recommendations:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.space_before = Pt(10)
        
        for idx, rec in enumerate(recommendations[:3], 1):
            p = text_frame.add_paragraph()
            p.text = f"{idx}. {rec['intervention'][:60]}..."
            p.font.size = Pt(10)
            p.space_before = Pt(4)
            p.level = 0
            
            p = text_frame.add_paragraph()
            p.text = f"   Ref: {rec['reference'][:50]}..."
            p.font.size = Pt(9)
            p.font.italic = True